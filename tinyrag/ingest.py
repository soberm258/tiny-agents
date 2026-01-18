import os
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Union

from loguru import logger

from tinyrag.utils import make_doc_id


def _read_text_file(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "gbk"):
        try:
            return path.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_md_file_to_text(path: Path) -> str:
    raw = _read_text_file(path)
    try:
        import markdown
        from bs4 import BeautifulSoup

        html_content = markdown.markdown(raw)
        soup = BeautifulSoup(html_content, "html.parser")
        return soup.get_text()
    except Exception:
        return raw


def _read_pdf_pages(path: Path) -> List[str]:
    import fitz  # PyMuPDF

    pdf_doc: fitz.Document = fitz.open(str(path))
    pages: List[str] = []
    for page in pdf_doc:
        text = page.get_text("text") or ""
        text = text.replace("\r\n", "\n")
        pages.append(text)
    return pages


def _read_docx_to_text(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    lines = [p.text for p in doc.paragraphs if p.text]
    return "\n".join(lines)


def _read_pptx_to_text(path: Path) -> str:
    from pptx import Presentation

    ppt = Presentation(str(path))
    lines: List[str] = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def _extract_texts_from_json_obj(obj: Any, *, text_key: str) -> List[str]:
    if obj is None:
        return []
    if isinstance(obj, str):
        return [obj]
    if isinstance(obj, list):
        out: List[str] = []
        for item in obj:
            out.extend(_extract_texts_from_json_obj(item, text_key=text_key))
        return out
    if isinstance(obj, dict):
        val = obj.get(text_key)
        if isinstance(val, str):
            return [val]
    return []


def load_docs_for_build(
    input_path: Union[str, Path],
    *,
    json_text_key: str = "completion",
    recursive: bool = True,
    suffix_allowlist: Optional[Iterable[str]] = None,
) -> List[Dict[str, Any]]:
    """
    将输入（文件或目录）读取为 build 可用的文档列表。
    返回值为 List[chunk]，每个元素是：
      {"id": str, "text": str, "meta": {...}}
    其中 meta 至少包含 source_path，PDF 还会包含 page。
    """
    input_path = Path(str(input_path))
    if not input_path.exists():
        raise FileNotFoundError(f"输入路径不存在：{input_path}")

    allow = None
    if suffix_allowlist is not None:
        allow = {s.lower().lstrip(".") for s in suffix_allowlist}

    def allow_suffix(path: Path) -> bool:
        if allow is None:
            return True
        return path.suffix.lower().lstrip(".") in allow

    files: List[Path] = []
    if input_path.is_dir():
        it = input_path.rglob("*") if recursive else input_path.glob("*")
        for p in it:
            if p.is_file() and allow_suffix(p):
                files.append(p)
    else:
        files = [input_path]

    docs: List[Dict[str, Any]] = []
    for file_path in files:
        suffix = file_path.suffix.lower().lstrip(".")
        if not suffix:
            continue

        try:
            if suffix == "pdf":
                pages = _read_pdf_pages(file_path)
                for idx, t in enumerate(pages, start=1):
                    meta = {"source_path": str(file_path), "page": idx, "type": "pdf"}
                    doc_id = make_doc_id(source_path=str(file_path), page=idx, record_index=0)
                    meta["doc_id"] = doc_id
                    docs.append({"id": doc_id, "text": t, "meta": meta})
            elif suffix == "txt":
                text = _read_text_file(file_path)
                if text.strip():
                    meta = {"source_path": str(file_path), "type": "txt"}
                    doc_id = make_doc_id(source_path=str(file_path), page=0, record_index=0)
                    meta["doc_id"] = doc_id
                    docs.append({"id": doc_id, "text": text, "meta": meta})
            elif suffix == "md":
                text = _read_md_file_to_text(file_path)
                if text.strip():
                    meta = {"source_path": str(file_path), "type": "md"}
                    doc_id = make_doc_id(source_path=str(file_path), page=0, record_index=0)
                    meta["doc_id"] = doc_id
                    docs.append({"id": doc_id, "text": text, "meta": meta})
            elif suffix == "docx":
                text = _read_docx_to_text(file_path)
                if text.strip():
                    meta = {"source_path": str(file_path), "type": "docx"}
                    doc_id = make_doc_id(source_path=str(file_path), page=0, record_index=0)
                    meta["doc_id"] = doc_id
                    docs.append({"id": doc_id, "text": text, "meta": meta})
            elif suffix == "pptx":
                text = _read_pptx_to_text(file_path)
                if text.strip():
                    meta = {"source_path": str(file_path), "type": "pptx"}
                    doc_id = make_doc_id(source_path=str(file_path), page=0, record_index=0)
                    meta["doc_id"] = doc_id
                    docs.append({"id": doc_id, "text": text, "meta": meta})
            elif suffix in ("json", "jsonl"):
                import json

                if suffix == "json":
                    obj = json.loads(_read_text_file(file_path))
                else:
                    obj = []
                    with open(file_path, "r", encoding="utf-8") as f:
                        for line in f:
                            line = line.strip()
                            if line:
                                obj.append(json.loads(line))
                texts = _extract_texts_from_json_obj(obj, text_key=json_text_key)
                for idx, t in enumerate([x for x in texts if x and x.strip()]):
                    meta = {"source_path": str(file_path), "type": suffix, "record_index": idx, "text_key": json_text_key}
                    doc_id = make_doc_id(source_path=str(file_path), page=0, record_index=idx)
                    meta["doc_id"] = doc_id
                    docs.append({"id": doc_id, "text": t, "meta": meta})
            else:
                logger.warning("不支持的文件类型，已跳过：{}", str(file_path))
        except Exception as e:
            logger.error("读取失败：{}，错误：{}", str(file_path), str(e))
            raise

    return docs
